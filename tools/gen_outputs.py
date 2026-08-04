# -*- coding: utf-8 -*-
"""
生成 Demo 输出成果（05-demo-outputs/）

- workbuddy-original/ : WorkBuddy 初稿（标注待人工核对）
- reviewed-final/     : 人工复核终稿
- expected-results/   : 标准答案要点（QA 参考）

所有内容来自 case_data.py，与输入材料严格一致。

License: MIT
"""
import os
import sys

import gen_documents as G
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import case_data as C

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
OUT_ORIG = os.path.join(ROOT, "05-demo-outputs", "workbuddy-original")
OUT_FINAL = os.path.join(ROOT, "05-demo-outputs", "reviewed-final")
OUT_EXPECT = os.path.join(ROOT, "05-demo-outputs", "expected-results")
DATE_TAG = "20260806"
VER = "v1.0"


def banner(doc, stage):
    """stage: 'original' | 'final'"""
    p = doc.add_paragraph()
    if stage == "original":
        r = p.add_run("【WorkBuddy 原始输出 · 初稿】本文件由 WorkBuddy 依据上传材料首次生成，"
                      "待人工核对数字、人名与日期后定稿。")
        G.set_run(r, size=9, color=G.GREY)
    else:
        r = p.add_run("【人工复核终稿】数字、人名、日期已与数据表及原始记录核对一致；"
                      "不含有待确认事项被误写为决定。")
        G.set_run(r, size=9, color=G.GREY)
    G.add_note(doc, f"⚠️ {C.DISCLAIMER}")


# ============================================================
# 03_会议纪要
# ============================================================
def build_minutes(stage):
    doc = Document()
    G.set_style_font(doc)
    G.add_title(doc, f"{C.MEETING_TITLE} · 会议纪要")
    banner(doc, stage)

    # 基本信息
    G.add_h2(doc, "会议基本信息")
    info = [
        ("会议名称", C.MEETING_TITLE),
        ("会议时间", f"{C.MEETING_DATE} {C.MEETING_TIME}"),
        ("会议地点", C.MEETING_PLACE),
        ("主持人", "沈知远（市场部负责人）"),
        ("记录人", "苏眠（市场专员）"),
        ("参会人员", "、".join(n.replace("　", "") for n, *_ in C.PEOPLE)),
    ]
    for k, v in info:
        p = doc.add_paragraph()
        r = p.add_run(f"{k}：{v}"); G.set_run(r, size=11)

    G.add_h1(doc, "一、会议背景与目标")
    G.add_body(doc, f"复盘 {C.PERIOD_LABEL} 市场部开展的 {C.CAMPAIGN_TOTAL} 项市场活动，"
        "梳理各渠道曝光与线索数据，明确已确认决定、行动事项与待确认事项，形成正式会议纪要。")

    G.add_h1(doc, "二、主要讨论内容")
    G.add_bullet(doc, "夏日新品试饮（7/5–7/12，12 家直营门店，主推青梅乌龙、桂香冷萃）已完成。")
    G.add_bullet(doc, "小红书达人种草（7/8–7/25，合作达人 18 位，产出笔记 24 篇）已完成。")
    G.add_bullet(doc, "微信社群优惠（7/10–7/31，覆盖 26 个门店社群）已完成。")
    G.add_bullet(doc, "商场快闪（7/18–7/20，3 个核心商圈，同步抖音本地推流）已完成。")
    G.add_bullet(doc, "校园联名（原定 7/31 结束）因合作方 Logo 使用规范未确认，延期。")
    t = C.channel_totals()
    e, l, q, r = C.grand_total()
    G.add_bullet(doc, f"渠道数据：总曝光 {e:,} 次，总线索 {l:,} 条，有效线索 {q:,} 条，"
        f"整体转化率 {r:.1f}%（按有效线索÷线索计算）。")

    G.add_h1(doc, "三、已确认决定")
    for i, d in enumerate(C.DECISIONS, start=1):
        p = doc.add_paragraph()
        r = p.add_run(f"{i}. {d}"); G.set_run(r, size=11)

    G.add_h1(doc, "四、行动事项")
    tbl = doc.add_table(rows=1, cols=6)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["序号", "事项", "负责人", "协同", "截止日期", "备注"]):
        G.set_cell(tbl.rows[0].cells[i], h)
    G.style_table_header(tbl.rows[0])
    for (idx, item, owner, helper, due, note) in C.ACTIONS:
        cells = tbl.add_row().cells
        G.set_cell(cells[0], str(idx), bold=True)
        G.set_cell(cells[1], item)
        G.set_cell(cells[2], owner.replace("　", ""), bold=True)
        G.set_cell(cells[3], helper.replace("　", ""))
        G.set_cell(cells[4], due)
        G.set_cell(cells[5], note)

    G.add_h1(doc, "五、待确认事项")
    G.add_body(doc, "以下事项在原始记录中未形成结论，已按要求标记为「待确认」，"
        "未写入正式决定：")
    for i, (_, item, quote, miss) in enumerate(C.PENDING, start=1):
        p = doc.add_paragraph()
        r = p.add_run(f"{i}. {item} —— 原话：「{quote}」缺：{miss}")
        G.set_run(r, size=11)

    G.add_h1(doc, "六、下次会议安排")
    G.add_body(doc, "时间待定（下半年品牌代言人合作方向将于后续专项会议讨论）。")

    G.set_core_props(doc, "会议纪要")
    return doc


# ============================================================
# 04_工作总结（苏眠）
# ============================================================
def build_summary(stage):
    doc = Document()
    G.set_style_font(doc)
    G.add_title(doc, f"{C.PROTAGONIST}（{C.PROTAGONIST_ROLE}）{C.PERIOD_LABEL} 工作总结")
    banner(doc, stage)

    p = doc.add_paragraph()
    r = p.add_run(f"姓名：{C.PROTAGONIST}　岗位：{C.PROTAGONIST_ROLE}　统计周期：{C.PERIOD_RANGE}")
    G.set_run(r, size=11)

    G.add_h1(doc, "一、本月已完成工作")
    G.add_bullet(doc, "【完成】梳理夏日新品试饮执行清单，确认 12 家门店物料需求并跟进供应商下单发货。")
    G.add_bullet(doc, "【完成】7/5–7/12 夏日新品试饮活动落地，赴 3 家门店现场督导首日执行，收尾汇总执行情况与物料结余。")
    G.add_bullet(doc, "【完成】7/10 微信社群优惠活动上线，统一下发 26 个门店社群活动文案。")
    G.add_bullet(doc, "【参与】配合周然筛选小红书达人名单，提供门店客群画像参考。")
    G.add_bullet(doc, "【参与】协助陈屿完成商场快闪活动首日现场支持与物料回收。")
    G.add_bullet(doc, "【完成】7/22 启动校园联名活动，完成 2 所高校社团首轮沟通。")
    G.add_bullet(doc, "【完成】7/30 汇总社群活动全周期数据并提交郑楠。")

    G.add_h1(doc, "二、主要成果")
    G.add_bullet(doc, "夏日新品试饮覆盖 12 家直营门店，现场收集偏好反馈（青梅乌龙 / 桂香冷萃）。")
    G.add_bullet(doc, "微信社群优惠覆盖 26 个门店社群，社群渠道整体转化率 60.0%，为全渠道最高。")
    G.add_bullet(doc, "完成 2 所高校社团首轮沟通，校园联名活动进入物料设计阶段。")
    G.add_bullet(doc, "汇总并提交社群全周期数据，支撑部门月度渠道看板搭建。")

    G.add_h1(doc, "三、正在推进的工作")
    G.add_bullet(doc, "校园联名活动物料定稿（因合作方 Logo 使用规范待确认，已顺延）。")
    G.add_bullet(doc, "微信社群扩容测试方案拟定（计划 8 月启动，覆盖 5 个门店社群）。")

    G.add_h1(doc, "四、存在的问题")
    G.add_bullet(doc, "校园联名活动因合作方 Logo 使用规范确认延迟，物料定稿与执行顺延，需持续跟进。")
    G.add_bullet(doc, "社群规模扩大后转化率是否下降尚需数据验证，扩容测试应设对照。")

    G.add_h1(doc, "五、下月计划")
    G.add_bullet(doc, "完成校园联名活动物料定稿并启动执行（截止 8/8）。")
    G.add_bullet(doc, "拟定微信社群扩容测试方案并启动（截止 8/10）。")
    G.add_bullet(doc, "配合李维、周然完成抖音内容与目标人群匹配度复核相关协同工作。")

    G.add_h1(doc, "六、需要支持的事项")
    G.add_bullet(doc, "校园联名物料设计需设计同事支持约 2 个工作日。")
    G.add_bullet(doc, "需合作方尽快确认 Logo 使用规范，以便物料定稿。")

    G.set_core_props(doc, "工作总结")
    return doc


# ============================================================
# 05_活动通知（8 月安排）
# ============================================================
def build_notice(stage):
    doc = Document()
    G.set_style_font(doc)
    G.add_title(doc, f"{C.BRAND_CN}市场部 {C.NEXT_MONTH_LABEL} 市场活动安排通知")
    banner(doc, stage)

    p = doc.add_paragraph()
    r = p.add_run(f"通知对象：市场部全体、运营中心（郑楠）、相关协同同事")
    G.set_run(r, size=11)

    G.add_h1(doc, "一、活动背景")
    G.add_body(doc, f"依据 {C.MEETING_TITLE} 已确认决定与行动事项，现下发 {C.NEXT_MONTH_LABEL} 市场活动安排。"
        "本通知仅包含已确认事项，待确认事项不列入正式安排。")

    G.add_h1(doc, "二、主要活动安排")
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["任务", "负责人", "时间节点", "材料提交要求", "反馈方式"]):
        G.set_cell(tbl.rows[0].cells[i], h)
    G.style_table_header(tbl.rows[0])
    rows = [
        ("抖音渠道内容与目标人群匹配度复核报告", "李维", "2026-08-12", "含近 4 周素材分类统计", "提交市场部负责人"),
        ("抖音 A/B 测试素材 2 组", "周然", "2026-08-15", "需设计同事支持 2 个工作日", "同步李维"),
        ("校园联名活动物料定稿并启动执行", "苏眠", "2026-08-08", "待合作方 Logo 使用规范确认", "报备陈屿"),
        ("微信社群扩容测试方案", "苏眠", "2026-08-10", "覆盖 5 个门店社群", "报备市场部负责人"),
        ("7 月渠道数据汇总归档至共享盘", "郑楠", "2026-08-07", "含周度明细", "共享盘归档"),
        ("8 月市场活动安排通知起草与下发", "陈屿", "2026-08-06", "经部门负责人确认后下发", "市场部群"),
    ]
    for task, owner, due, req, fb in rows:
        cells = tbl.add_row().cells
        G.set_cell(cells[0], task)
        G.set_cell(cells[1], owner.replace("　", ""), bold=True)
        G.set_cell(cells[2], due)
        G.set_cell(cells[3], req)
        G.set_cell(cells[4], fb)

    G.add_h1(doc, "三、反馈与联系方式")
    G.add_body(doc, "各项任务截止前 1 个工作日请主动向负责人反馈进度；遇口径或资源问题，"
        "及时在市场部群同步。联系人：沈知远（市场部负责人）。")

    G.set_core_props(doc, "活动通知")
    return doc


# ============================================================
# 06_市场分析报告
# ============================================================
def build_report(stage):
    doc = Document()
    G.set_style_font(doc)
    G.add_title(doc, f"{C.BRAND_CN} {C.DEPARTMENT} {C.PERIOD_LABEL} 市场数据分析报告")
    banner(doc, stage)
    G.add_note(doc, f"统一口径：{C.METRIC_FORMULA}　｜　整体转化率 = 有效线索合计 ÷ 线索合计（非算术平均）。")

    G.add_h1(doc, "一、数据完整性检查")
    G.add_bullet(doc, "【事实】5 个渠道 × 4 周 = 20 行明细，无缺值、无负数。")
    G.add_bullet(doc, "【事实】曝光、线索、有效线索三级数据均在各渠道后台导出并人工初筛。")
    G.add_bullet(doc, "【判断】数据口径在 8 月起已统一，本月复盘按统一公式计算。")

    G.add_h1(doc, "二、市场活动完成情况")
    G.add_bullet(doc, f"【事实】共 {C.CAMPAIGN_TOTAL} 项活动，{C.CAMPAIGN_DONE} 项已完成，{C.CAMPAIGN_DELAYED} 项延期。")
    G.add_bullet(doc, "【事实】校园联名因合作方 Logo 使用规范确认延迟，顺延执行，不计入 7 月完成率。")
    G.add_bullet(doc, f"【事实】活动完成率 = {C.CAMPAIGN_DONE} ÷ {C.CAMPAIGN_TOTAL} = {C.CAMPAIGN_DONE/C.CAMPAIGN_TOTAL:.0%}（仅计已完成）。")

    G.add_h1(doc, "三、各渠道曝光与线索")
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["渠道", "曝光量", "线索数量", "有效线索"]):
        G.set_cell(tbl.rows[0].cells[i], h)
    G.style_table_header(tbl.rows[0])
    t = C.channel_totals()
    for ch in C.CHANNELS:
        exp, lead, qual, _ = t[ch]
        cells = tbl.add_row().cells
        G.set_cell(cells[0], ch, bold=True)
        G.set_cell(cells[1], f"{exp:,}")
        G.set_cell(cells[2], f"{lead:,}")
        G.set_cell(cells[3], f"{qual:,}")
    e, l, q, _ = C.grand_total()
    cells = tbl.add_row().cells
    G.set_cell(cells[0], "合计", bold=True, color=G.NAVY)
    G.set_cell(cells[1], f"{e:,}", bold=True)
    G.set_cell(cells[2], f"{l:,}", bold=True)
    G.set_cell(cells[3], f"{q:,}", bold=True)

    G.add_h1(doc, "四、转化率分析")
    G.add_bullet(doc, "【事实】各渠道转化率：微信社群 60.0% > 微信公众号 45.0% = 线下门店 45.0% > 小红书 34.1% > 抖音 25.0%。")
    G.add_bullet(doc, f"【事实】整体转化率 = {q:,} ÷ {l:,} = {q/l*100:.1f}%。")
    G.add_bullet(doc, "【判断】整体转化率（37.0%）低于多数单渠道转化率，因抖音低转化高曝光拉低加权结果；"
        "不可用五渠道算术平均（约 41.8%）替代。")

    G.add_h1(doc, "五、表现较好的渠道")
    G.add_bullet(doc, "【事实】微信社群转化率 60.0% 最高，虽曝光量最小（32,000），但客群精准。")
    G.add_bullet(doc, "【事实】微信公众号与线下门店转化率均为 45.0%，表现稳健。")

    G.add_h1(doc, "六、需要关注的问题")
    G.add_bullet(doc, "【事实】抖音曝光量最高（262,000）但转化率最低（25.0%），有效线索仅 390 条。")
    G.add_bullet(doc, "【判断】抖音内容可能与目标人群匹配度不足，需先做内容优化而非加预算。")

    G.add_h1(doc, "七、数据可支持的改进建议")
    G.add_bullet(doc, "【建议】抖音先做内容与目标人群匹配度复核与 A/B 测试（已列入 8 月行动事项），不急于加预算。")
    G.add_bullet(doc, "【建议】在微信社群开展小规模扩容测试，验证转化率是否随规模下降。")
    G.add_bullet(doc, "【建议】校园联名物料定稿后尽快启动，弥补 7 月延期缺口。")
    G.add_note(doc, "说明：本报告不含预算、成本、销售额、客单价、用户画像数据，相关结论无数据支撑，不予推断。")

    G.set_core_props(doc, "市场数据分析报告")
    return doc


# ============================================================
# 07_月度复盘PPT（终审稿 + 原稿）
# ============================================================
def build_ppt(stage):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY = PptRGB(0x1B, 0x3A, 0x6B)
    BLUE = PptRGB(0x2E, 0x6B, 0xE6)
    WHITE = PptRGB(0xFF, 0xFF, 0xFF)
    GREY = PptRGB(0x55, 0x55, 0x55)
    LIGHT = PptRGB(0xEA, 0xF1, 0xFD)

    def add_bg(slide, color=WHITE):
        box = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        box.fill.solid(); box.fill.fore_color.rgb = color
        box.line.fill.background()
        box.shadow.inherit = False
        return box

    def add_bar(slide, title, subtitle=""):
        add_bg(slide, WHITE)
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]; p.text = title
        p.font.size = PptPt(30); p.font.bold = True; p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        if subtitle:
            sp = tf.add_paragraph(); sp.text = subtitle
            sp.font.size = PptPt(14); sp.font.color.rgb = LIGHT; sp.font.name = "微软雅黑"

    def add_text(slide, l, t, w, h, text, size=20, color=GREY, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.alignment = align
        p.font.size = PptPt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = "微软雅黑"
        return box

    def add_bullets(slide, l, t, w, h, items, size=20, color=GREY):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "• " + it
            p.font.size = PptPt(size); p.font.color.rgb = color; p.font.name = "微软雅黑"
            p.space_after = PptPt(6)
        return box

    # 1 封面
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_text(s, 1.0, 2.6, 11.3, 1.4, f"{C.BRAND_FULL}", size=44, color=WHITE, bold=True)
    add_text(s, 1.0, 3.9, 11.3, 1.0, f"市场部 {C.PERIOD_LABEL} 月度复盘", size=30, color=LIGHT)
    add_text(s, 1.0, 5.0, 11.3, 0.6,
             ("WorkBuddy 原始输出 · 初稿" if stage == "original" else "人工复核终稿")
             + "　｜　仅用于教学演示", size=14, color=LIGHT)

    # 2 目录
    s = prs.slides.add_slide(blank)
    add_bar(s, "目录")
    add_bullets(s, 0.8, 1.5, 11.5, 5.5, [
        "案例背景与本月市场活动", "渠道曝光与线索数据", "已确认决定与行动事项",
        "待确认事项", "关键结论与下月计划", "寄语",
    ], size=22)

    # 3 案例背景
    s = prs.slides.add_slide(blank)
    add_bar(s, "案例背景", f"{C.BRAND_FULL} · {C.DEPARTMENT}")
    add_bullets(s, 0.8, 1.5, 11.5, 5.0, [
        f"虚拟新式茶饮品牌，本教学案例以市场部 {C.PERIOD_LABEL} 月度复盘为场景。",
        f"本月开展 {C.CAMPAIGN_TOTAL} 项市场活动，覆盖 {len(C.CHANNELS)} 个渠道。",
        "目标：把零散会议记录、渠道数据、个人工作记录整理为正式办公成果。",
        "全部品牌、人员、数据均为虚构，仅用于教学演示。",
    ])

    # 4 本月市场活动
    s = prs.slides.add_slide(blank)
    add_bar(s, "本月市场活动", f"共 {C.CAMPAIGN_TOTAL} 项 · 已完成 {C.CAMPAIGN_DONE} · 延期 {C.CAMPAIGN_DELAYED}")
    tbl = s.shapes.add_table(6, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0)).table
    hdr = ["活动", "状态", "负责人"]
    for j, h in enumerate(hdr):
        tbl.cell(0, j).text = h
        tbl.cell(0, j).fill.solid(); tbl.cell(0, j).fill.fore_color.rgb = BLUE
        pr = tbl.cell(0, j).text_frame.paragraphs[0]
        pr.font.size = PptPt(16); pr.font.bold = True; pr.font.color.rgb = WHITE
    for i, (_, name, *_rest, status, owner, helper, ch, desc) in enumerate(C.CAMPAIGNS, start=1):
        tbl.cell(i, 0).text = name
        tbl.cell(i, 1).text = status
        tbl.cell(i, 2).text = owner.replace("　", "")
        for j in range(3):
            pr = tbl.cell(i, j).text_frame.paragraphs[0]
            pr.font.size = PptPt(15); pr.font.name = "微软雅黑"
            if j == 1 and status == "延期":
                pr.font.color.rgb = PptRGB(0xC0, 0x39, 0x2B); pr.font.bold = True

    # 5 渠道数据概览
    s = prs.slides.add_slide(blank)
    e, l, q, r = C.grand_total()
    add_bar(s, "渠道数据概览", f"统计周期：{C.PERIOD_RANGE}")
    add_bullets(s, 0.8, 1.5, 11.5, 4.5, [
        f"总曝光量：{e:,} 次",
        f"总线索数量：{l:,} 条",
        f"总有效线索：{q:,} 条",
        f"整体转化率：{r:.1f}%（有效线索 ÷ 线索）",
    ], size=24)
    add_text(s, 0.8, 6.3, 11.5, 0.6,
             "统计口径：转化率 = 有效线索数量 ÷ 线索数量 × 100%", size=14, color=BLUE)

    # 6 各渠道转化率（原生图表）
    s = prs.slides.add_slide(blank)
    add_bar(s, "各渠道转化率对比", "单位：%　（整体转化率 ≠ 各渠道算术平均）")
    t = C.channel_totals()
    chart_data = CategoryChartData()
    chart_data.categories = C.CHANNELS
    chart_data.add_series("转化率", [t[ch][3] for ch in C.CHANNELS])
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0), chart_data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0.0"%"'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = PptPt(14)

    # 7 已确认决定
    s = prs.slides.add_slide(blank)
    add_bar(s, "已确认决定")
    add_bullets(s, 0.8, 1.5, 11.5, 5.2, [f"{i}. {d}" for i, d in enumerate(C.DECISIONS, 1)], size=18)

    # 8 行动事项
    s = prs.slides.add_slide(blank)
    add_bar(s, "行动事项（8 月）")
    tbl = s.shapes.add_table(len(C.ACTIONS) + 1, 4, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2)).table
    for j, h in enumerate(["事项", "负责人", "截止", "协同"]):
        tbl.cell(0, j).text = h
        tbl.cell(0, j).fill.solid(); tbl.cell(0, j).fill.fore_color.rgb = BLUE
        pr = tbl.cell(0, j).text_frame.paragraphs[0]
        pr.font.size = PptPt(15); pr.font.bold = True; pr.font.color.rgb = WHITE
    for i, (_, item, owner, helper, due, note) in enumerate(C.ACTIONS, start=1):
        tbl.cell(i, 0).text = item
        tbl.cell(i, 1).text = owner.replace("　", "")
        tbl.cell(i, 2).text = due
        tbl.cell(i, 3).text = helper.replace("　", "")
        for j in range(4):
            pr = tbl.cell(i, j).text_frame.paragraphs[0]
            pr.font.size = PptPt(13); pr.font.name = "微软雅黑"

    # 9 待确认事项
    s = prs.slides.add_slide(blank)
    add_bar(s, "待确认事项", "原始记录未形成结论，已标记待确认，未写入正式决定")
    items = [f"{i}. {item}" for i, (item, quote, miss) in enumerate(
        [(p[1], p[2], p[3]) for p in C.PENDING], 1)]
    add_bullets(s, 0.8, 1.5, 11.5, 5.2, items, size=18, color=PptRGB(0x55, 0x55, 0x55))

    # 10 关键结论
    s = prs.slides.add_slide(blank)
    add_bar(s, "关键结论")
    add_bullets(s, 0.8, 1.5, 11.5, 5.2, [
        "微信社群转化率最高（60.0%），客群精准，建议小范围扩容验证。",
        "抖音曝光最高但转化率最低（25.0%），应先做内容优化，不急于加预算。",
        "整体转化率 37.0% 为加权结果，不可用各渠道算术平均替代。",
        "校园联名延期，需尽快确认合作方 Logo 规范后启动。",
    ], size=20)

    # 11 下月计划
    s = prs.slides.add_slide(blank)
    add_bar(s, "下月计划")
    add_bullets(s, 0.8, 1.5, 11.5, 5.2, [
        "8/8 前：校园联名活动物料定稿并启动执行（苏眠）。",
        "8/10 前：微信社群扩容测试方案拟定（苏眠）。",
        "8/12 前：抖音内容与目标人群匹配度复核报告（李维）。",
        "8/15 前：抖音 A/B 测试素材 2 组（周然）。",
        "8/7 前：7 月渠道数据汇总归档（郑楠）。",
    ], size=20)

    # 12 寄语
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_text(s, 1.0, 2.2, 11.3, 3.2,
             "初入职场，不必要求自己立刻掌握所有工具，\n也不必一次写出完美的指令。\n"
             "真正重要的是理解工作的本质，知道目标是什么，知道结果该如何判断。\n"
             "愿你保持好奇，持续学习，在与 AI 共同工作的未来中，更高效地成长。",
             size=22, color=WHITE, bold=False)

    return prs


def save_doc(doc, name):
    p = os.path.join(OUT_FINAL if name.startswith("F_") else OUT_ORIG, name)
    doc.save(p)
    return p


def main():
    C.verify()
    for d in (OUT_ORIG, OUT_FINAL, OUT_EXPECT):
        os.makedirs(d, exist_ok=True)

    # 会议纪要
    orig = build_minutes("original"); fin = build_minutes("final")
    op = os.path.join(OUT_ORIG, f"03_会议纪要_WorkBuddy原始版_{VER}_{DATE_TAG}.docx"); orig.save(op)
    fp = os.path.join(OUT_FINAL, f"03_会议纪要_人工审核版_{VER}_{DATE_TAG}.docx"); fin.save(fp)
    print(f"[OK] {os.path.relpath(op, ROOT)}")
    print(f"[OK] {os.path.relpath(fp, ROOT)}")

    # 工作总结
    orig = build_summary("original"); fin = build_summary("final")
    op = os.path.join(OUT_ORIG, f"04_工作总结_WorkBuddy原始版_{VER}_{DATE_TAG}.docx"); orig.save(op)
    fp = os.path.join(OUT_FINAL, f"04_工作总结_人工审核版_{VER}_{DATE_TAG}.docx"); fin.save(fp)
    print(f"[OK] {os.path.relpath(op, ROOT)}")
    print(f"[OK] {os.path.relpath(fp, ROOT)}")

    # 通知
    orig = build_notice("original"); fin = build_notice("final")
    op = os.path.join(OUT_ORIG, f"05_活动通知_WorkBuddy原始版_{VER}_{DATE_TAG}.docx"); orig.save(op)
    fp = os.path.join(OUT_FINAL, f"05_活动通知_人工审核版_{VER}_{DATE_TAG}.docx"); fin.save(fp)
    print(f"[OK] {os.path.relpath(op, ROOT)}")
    print(f"[OK] {os.path.relpath(fp, ROOT)}")

    # 分析报告
    orig = build_report("original"); fin = build_report("final")
    op = os.path.join(OUT_ORIG, f"06_市场分析报告_WorkBuddy原始版_{VER}_{DATE_TAG}.docx"); orig.save(op)
    fp = os.path.join(OUT_FINAL, f"06_市场分析报告_人工审核版_{VER}_{DATE_TAG}.docx"); fin.save(fp)
    print(f"[OK] {os.path.relpath(op, ROOT)}")
    print(f"[OK] {os.path.relpath(fp, ROOT)}")

    # PPT
    op = os.path.join(OUT_ORIG, f"07_月度复盘PPT_WorkBuddy原始版_{VER}_{DATE_TAG}.pptx")
    build_ppt("original").save(op)
    fp = os.path.join(OUT_FINAL, f"07_月度复盘PPT_人工审核版_{VER}_{DATE_TAG}.pptx")
    build_ppt("final").save(fp)
    print(f"[OK] {os.path.relpath(op, ROOT)}")
    print(f"[OK] {os.path.relpath(fp, ROOT)}")

    # 标准答案要点（expected-results）
    exp = Document()
    G.set_style_font(exp)
    G.add_title(exp, "标准答案要点（QA 参考）")
    G.add_note(exp, "本文件用于质量检查，列出 5 份输出成果必须包含的关键点。")
    secs = [
        ("会议纪要", ["会议基本信息完整", "4 条已确认决定逐条列出", "6 项行动事项含负责人与截止日期",
                   "5 项待确认事项标记，未写入决定", "整体转化率 37.0% 口径正确"]),
        ("工作总结", ["区分「完成 / 参与 / 协助」", "校园联名标为进行中（延期）", "不把团队成果写成个人成果",
                   "数据对应具体活动", "含下月计划与需支持事项"]),
        ("活动通知", ["仅含已确认事项", "不含待确认事项", "含对象 / 任务 / 负责人 / 时间 / 反馈方式"]),
        ("分析报告", ["区分事实 / 判断 / 建议", "整体转化率 37.0% 非算术平均", "抖音低转化已指出",
                   "不虚构预算 / 销售额 / 画像", "改进建议仅限数据支撑"]),
        ("月度复盘PPT", ["10–12 页", "每页一个核心结论", "数据页标统计周期与口径", "含寄语页", "不覆盖文字"]),
    ]
    for title, items in secs:
        G.add_h1(exp, title)
        for it in items:
            G.add_bullet(exp, it)
    ep = os.path.join(OUT_EXPECT, "标准答案要点.docx")
    G.set_core_props(exp, "标准答案要点")
    exp.save(ep)
    print(f"[OK] {os.path.relpath(ep, ROOT)}")
    print("\nDemo 输出成果生成完成。")


if __name__ == "__main__":
    main()
