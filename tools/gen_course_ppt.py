# -*- coding: utf-8 -*-
"""
生成课程 PPT 与 PDF（02-slides/）

- workbuddy-beginner-training.pptx  （科技风，16:9，约 44 页）
- workbuddy-beginner-training.pdf   （与 PPT 内容一致，可打印）
- slide-outline.md                 （页面索引）
- slide-notes.md                   （讲师备注）

本版本为「功能导向」结构：
  先讲 WorkBuddy 各功能是什么 / 干什么 / 怎么用 / 在哪里；
  用统一案例（青屿茶研）演示每个功能的具体使用流程；
  进阶技巧：先给方案再执行、出错后怎么办；
  Demo 结果仅顺带展示。

不含真实界面截图（需学员运行 WorkBuddy 后自行采集）。流程图 / 界面示意
均由 python-pptx 绘制，直观展示功能位置与使用步骤。

License: MIT
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas as RLCanvas

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
SLIDE_DIR = os.path.join(ROOT, "02-slides")
CHART_DIR = os.path.join(ROOT, "08-design-assets", "charts")

# ── 主题色 ──
NAVY = PptRGB(0x1B, 0x3A, 0x6B)
BLUE = PptRGB(0x2E, 0x6B, 0xE6)
CYAN = PptRGB(0x17, 0xC0, 0xC7)
LIGHT = PptRGB(0xEA, 0xF1, 0xFD)
WHITE = PptRGB(0xFF, 0xFF, 0xFF)
GREY = PptRGB(0x55, 0x55, 0x55)
RED = PptRGB(0xC0, 0x39, 0x2B)
GREEN = PptRGB(0x27, 0xAE, 0x60)
INK = PptRGB(0x22, 0x2A, 0x35)
AMBER = PptRGB(0xE6, 0x7E, 0x22)

NAVY_H = (0x1B, 0x3A, 0x6B)
BLUE_H = (0x2E, 0x6B, 0xE6)
LIGHT_H = (0xEA, 0xF1, 0xFD)
GREY_H = (0x55, 0x55, 0x55)
RED_H = (0xC0, 0x39, 0x2B)
GREEN_H = (0x27, 0xAE, 0x60)


# ============================================================
# 幻灯片内容（功能导向，单一来源）
# ============================================================
SLIDES = [
    # ═══════════ 封面 ═══════════
    {"type": "cover", "kicker": "WorkBuddy 零基础办公任务实战教程",
     "title": "认识 WorkBuddy：功能 · 用法 · 实战",
     "bullets": ["先搞懂每个功能是什么、干什么、怎么用、在哪里",
                 "用统一案例演示各功能的真实使用流程",
                 "再学进阶技巧：先给方案再执行、出错怎么办"],
     "note": "开场点题：今天不堆砌案例，而是把 WorkBuddy 的功能地图讲清楚，再用一个案例把流程串起来。"},

    # ═══════════ 课程目标 ═══════════
    {"type": "content", "chapter": "课程目标", "kicker": "学完能做什么",
     "title": "本课程的三条目标",
     "bullets": ["能说出 WorkBuddy 的核心功能分别解决哪类办公问题",
                 "能在界面上找到对应入口，并跑通一个完整任务流程",
                 "掌握两个进阶习惯：先要方案再执行、出错了按流程纠正"],
     "note": "强调目标偏「认知 + 操作」，不是背按钮。"},

    # ═══════════ 第一章 ═══════════
    {"type": "divider", "kicker": "第一章", "title": "为什么学习 WorkBuddy",
     "note": "先建立动机，再讲功能。"},
    {"type": "content", "chapter": "第一章 为什么学习 WorkBuddy", "kicker": "核心信息",
     "title": "本质、效率与未来",
     "bullets": ["本质：WorkBuddy 不是替你思考，而是把目标、材料、约束变成可执行任务",
                 "效率：减少澄清、整理、排版、转换中的机械劳动",
                 "未来：能定义目标、组织材料、检查结果，并与 AI 协同"],
     "note": "三句话分别回开场、中段、结语，须反复强化。"},
    {"type": "content", "chapter": "第一章 为什么学习 WorkBuddy", "kicker": "痛点",
     "title": "新员工的困境，正是功能要解决的",
     "bullets": ["说不清需求 → 用「一句话下达任务 + 先提问」功能澄清",
                 "材料散乱 → 用「添加上下文 / 选择工作空间」功能集中输入",
                 "反复排版 → 用「文档 / 表格 / PPT 生成」功能一次产出",
                 "不敢交付 → 用「结果区验收 + 四查」功能把控质量"],
     "note": "把困境和后面的功能一一对应，让学员知道学这些有什么用。"},

    # ═══════════ 第二章 认识界面与起步 ═══════════
    {"type": "divider", "kicker": "第二章", "title": "认识 WorkBuddy 界面与起步",
     "note": "先知道东西在哪、怎么开始，才能谈功能。"},
    {"type": "content", "chapter": "第二章 认识界面与起步", "kicker": "定义",
     "title": "WorkBuddy 是什么",
     "bullets": ["腾讯出品的全场景 AI 办公工作台：说一句话，它规划并执行任务、交付成果",
                 "能读你的本地文件、写文件、跨步骤推进，不只是问答",
                 "与传统 AI 对话的区别：能实际操作文件、有任务连续性、交付可验收结果"],
     "note": "用「能不能处理我的文件」来区分问答工具和 WorkBuddy。"},
    {"type": "mockup", "chapter": "第二章 认识界面与起步", "kicker": "界面",
     "title": "界面三区总览（东西在哪里）",
     "note": "整页都是「在哪里」。建议手指屏幕三区逐一讲。"},
    {"type": "content", "chapter": "第二章 认识界面与起步", "kicker": "起步①",
     "title": "第一步：确认积分充足",
     "bullets": ["设计、生成图片 / 视频等多模态功能会消耗积分（信用点）",
                 "开始重任务前，先在账号 / 设置中确认积分余额充足",
                 "积分不足时系统会提示，避免生成到一半中断"],
     "note": "用户特别提到「确认积分充足」作为流程起点，务必放在第一步。"},
    {"type": "content", "chapter": "第二章 认识界面与起步", "kicker": "起步②",
     "title": "第二步：选择工作模式",
     "bullets": ["Craft（你说我做）：直接动手产出结果，最快",
                 "Plan（先想后做）：先给方案与步骤，你确认后再执行",
                 "Ask（只看不说）：只回答、只读，不改动任何文件",
                 "在哪里：对话区顶部或新建任务时的模式切换"],
     "note": "新手建议先用 Plan，看清步骤再放手。"},
    {"type": "content", "chapter": "第二章 认识界面与起步", "kicker": "起步③",
     "title": "第三步：一句话下达任务 + 选工作空间",
     "bullets": ["在对话区底部输入框，用自然语言写需求，如「把这份会议记录整理成纪要」",
                 "点输入框左下角「选择工作空间」，指定本次任务读取 / 生成文件的目录",
                 "不选则用默认目录；指定后 WorkBuddy 只在该目录内读写",
                 "回车即创建任务，新任务出现在左侧任务列表"],
     "note": "强调工作空间决定文件范围，避免误读其他资料。"},
    {"type": "content", "chapter": "第二章 认识界面与起步", "kicker": "起步④",
     "title": "第四步：添加上下文的 4 种方式",
     "bullets": ["@ 引用：在输入中 @ 文件 / 文档 / 规则，精准指向材料",
                 "粘贴截图：Ctrl/Cmd+V 直接粘贴剪贴板图片",
                 "上传文件：点上传按钮或拖拽文件到输入框",
                 "补充说明：在描述里写明目标、输入、输出格式、约束、时限"],
     "note": "上下文越清楚，结果越准。这是后面所有 Demo 的公共前置。"},

    # ═══════════ 第三章 核心功能逐个看 ═══════════
    {"type": "divider", "kicker": "第三章", "title": "核心功能逐个看",
     "note": "本章是重点。每个功能讲清是什么/干什么/怎么用/在哪里，并用案例演示流程。"},
    {"type": "content", "chapter": "第三章 核心功能逐个看", "kicker": "地图",
     "title": "WorkBuddy 功能地图（10 大功能）",
     "bullets": ["① 文档生成  ② 表格与数据分析  ③ PPT 生成  ④ 设计 / 多模态",
                 "⑤ 深度研究与联网  ⑥ 批量文件处理  ⑦ 代码 / 工具执行",
                 "⑧ 连接器 / MCP（GitHub、邮箱…）  ⑨ 自动化（定时任务）  ⑩ 专家",
                 "下面逐一看定义、用法、入口，并用青屿茶研案例演示关键流程"],
     "note": "先给全景，学员心里有地图再逐个深入。"},

    # —— 功能① 文档生成 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ①",
     "title": "文档生成（纪要 / 报告 / 通知）",
     "what": "把零散材料整理成结构化的 Word 文档。",
     "do": "会议纪要、工作总结、活动通知、调研报告、README 等。",
     "how": ["上传原始材料 + 模板", "写明结构要求与约束", "先要大纲，确认后生成 .docx"],
     "where": "对话区直接下达任务；结果区「产物」下载 .docx。",
     "demo": "案例演示：会议原始记录 → 标准会议纪要（见后页流程）",
     "note": "文档生成是最常用的功能，先讲透。"},
    {"type": "flow", "chapter": "第三章 核心功能逐个看", "kicker": "流程 ①",
     "title": "文档生成 · Demo 使用流程",
     "steps": ["创建任务\n选工作空间", "上传会议\n记录+模板", "下达指令\n(只据材料)", "先出结构\n+待确认项", "核对人名/\n数字→导出"],
     "note": "每步都对应界面操作：新建任务→底部上传→输入框写约束→确认大纲→结果区下载。"},

    # —— 功能② 表格与数据分析 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ②",
     "title": "表格与数据分析",
     "what": "读取 Excel / CSV，做汇总、公式、对比与可视化。",
     "do": "渠道曝光汇总、转化率计算、趋势图、异常预警。",
     "how": ["上传数据表（含公式更佳）", "说明统计口径与指标定义", "要求图表 + 关键结论"],
     "where": "对话区下达任务；结果区看产物，或让它写回 .xlsx。",
     "demo": "案例演示：渠道数据 → 整体转化率（加权 37.0%）",
     "note": "强调「合计类指标要加权」，这是课程核心纠错点。"},
    {"type": "flow", "chapter": "第三章 核心功能逐个看", "kicker": "流程 ②",
     "title": "数据分析 · Demo 使用流程",
     "steps": ["上传渠道\n数据表", "明确口径\n(有效÷线索)", "生成各渠道\n汇总", "计算整体\n(加权)", "出图+结论\n(非平均)"],
     "note": "核心纠错点：整体转化率 = 总有效线索 ÷ 总线索 = 1744 ÷ 4710 = 37.0%（加权），不能用五渠道算术平均 ~41.8% 替代。"},

    # —— 功能③ PPT 生成 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ③",
     "title": "PPT 生成（报告 / 复盘）",
     "what": "从需求或已有文档，生成结构化的演示文稿。",
     "do": "月度复盘、项目汇报、培训课件、产品介绍。",
     "how": ["提供输入材料与目标页数", "先要逐页大纲并确认", "再生成 .pptx，逐页可改"],
     "where": "对话区下达任务；结果区「产物」打开 / 下载 .pptx。",
     "demo": "案例演示：纪要+总结+通知+分析 → 12 页复盘 PPT",
     "note": "PPT 是综合输出，放在多个单功能之后讲最顺。"},
    {"type": "flow", "chapter": "第三章 核心功能逐个看", "kicker": "流程 ③",
     "title": "PPT 生成 · Demo 使用流程",
     "steps": ["汇总已完成的\n4份文档", "下达「生成\n复盘PPT」", "先出逐页\n大纲", "确认/修改\n大纲", "生成12页\nPPT"],
     "note": "凸显「先大纲后生成」的进阶技巧。"},

    # —— 功能④ 设计 / 多模态 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ④",
     "title": "设计 / 多模态（图 / 海报 / Logo）",
     "what": "用文字描述生成图片、海报、Logo、配图等视觉素材。",
     "do": "活动海报、社交媒体配图、品牌标识、PPT 插图。",
     "how": ["用文字描述画面 / 风格 / 尺寸", "确认积分充足（消耗信用点）", "生成后选优或局部修改"],
     "where": "对话区下达「生成图片」类任务；结果区看成品。",
     "demo": "案例演示：校园联名活动 → 一张预热海报",
     "note": "提醒积分：多模态生成前有「确认积分充足」这一步。"},
    {"type": "flow", "chapter": "第三章 核心功能逐个看", "kicker": "流程 ④",
     "title": "设计功能 · Demo 使用流程",
     "steps": ["确认积分\n充足", "下达「生成\n活动海报」", "描述主题/\n风格/尺寸", "生成候选\n多张", "选定/微调\n导出"],
     "note": "设计功能消耗积分，所以紧跟在「确认积分充足」之后最自然。"},

    # —— 功能⑤ 深度研究 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑤",
     "title": "深度研究与联网",
     "what": "联网检索 + 多轮调研，输出带出处的报告。",
     "do": "行业趋势、竞品分析、资料查证、方案对比。",
     "how": ["提出研究问题与目标", "让它分步检索并标注来源", "要求结论 + 引用 + 不确定项"],
     "where": "对话区下达「调研…」任务；开启联网能力。",
     "demo": "可扩展：调研「新茶饮 7 月营销趋势」辅助决策",
     "note": "联网结果要核对出处，避免编造来源。"},

    # —— 功能⑥ 批量文件处理 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑥",
     "title": "批量文件处理",
     "what": "对一类文件批量整理、重命名、格式转换。",
     "do": "按日期分类图片、批量改名、docx↔pdf 转换。",
     "how": ["选择工作空间（目标目录）", "说明规则：按 X 分类 / 命名为 Y", "让它批量执行并汇报变更"],
     "where": "对话区下达任务；结果区「变更」查看改了哪些。",
     "demo": "可扩展：把活动素材按渠道文件夹自动归类",
     "note": "强调先看「变更」预览，确认无误再保留。"},

    # —— 功能⑦ 代码 / 工具执行 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑦",
     "title": "代码 / 工具执行",
     "what": "在受控环境运行脚本、处理数据、生成文件。",
     "do": "用 Python 批量生成表格 / 图表、跑数据处理。",
     "how": ["说明要生成的产物", "让它写脚本并执行", "在结果区查看运行产物"],
     "where": "对话区下达任务；WorkBuddy 在沙箱中执行。",
     "demo": "本教程所有图表即由脚本生成（tools/gen_charts.py）",
     "note": "说明：生成脚本也是 WorkBuddy 能做的事，体现「可执行」。"},

    # —— 功能⑧ 连接器 / MCP ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑧",
     "title": "连接器 / MCP（连接外部系统）",
     "what": "接入 GitHub、邮箱、腾讯文档等，让 WB 直接操作它们。",
     "do": "推送代码到 GitHub、读邮件、写在线文档。",
     "how": ["在连接器中心开启并授权", "下达任务时指明目标系统", "按提示完成授权确认"],
     "where": "左侧 / 设置中的「连接器」；授权后对话区直接调用。",
     "demo": "案例演示：把成品仓库推送到 GitHub（见后页流程）",
     "note": "连接器需授权，涉及外部系统操作要谨慎。"},
    {"type": "flow", "chapter": "第三章 核心功能逐个看", "kicker": "流程 ⑧",
     "title": "连接器 · Demo 使用流程（发布到 GitHub）",
     "steps": ["开启 GitHub\n连接器授权", "准备本地\n仓库文件", "下达「推送\n到 GitHub」", "确认仓库/\n分支", "推送成功\n获链接"],
     "note": "本教程仓库即由此流程发布，作为可复现范例。"},

    # —— 功能⑨ 自动化 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑨",
     "title": "自动化（定时 / 周期任务）",
     "what": "把重复任务设成定时或周期性自动运行。",
     "do": "每周数据汇总、每日舆情简报、定时提醒。",
     "how": ["描述任务与触发时间（如每周一 9 点）", "WorkBuddy 登记自动化", "到期自动执行并通知"],
     "where": "自动化入口（任务管理 / 设置）；按周期自动跑。",
     "demo": "可扩展：每周一自动汇总上周渠道数据",
     "note": "适合「每次都一样的重复活」，解放双手。"},

    # —— 功能⑩ 专家 ——
    {"type": "feature", "chapter": "第三章 核心功能逐个看", "kicker": "功能 ⑩",
     "title": "专家（领域智能体）",
     "what": "调用预先配置好的领域专家，处理专业问题。",
     "do": "财务分析、法律咨询、设计评审、行业研究。",
     "how": ["从专家中心选择对应专家", "像对话一样把问题交给它", "结合 WB 通用能力协同"],
     "where": "专家中心 / 专家入口；选定后进入专家对话。",
     "demo": "可扩展：用「财务专家」复核转化率口径",
     "note": "专家是「专业能力包」，和通用对话互补。"},

    # —— 结果区 / 侧边栏 ——
    {"type": "content", "chapter": "第三章 核心功能逐个看", "kicker": "验收",
     "title": "结果区：查看与验收产物",
     "bullets": ["产物：直接预览生成的文档 / 表格 / PPT / 图片",
                 "全部文件：本次任务产生的所有文件清单",
                 "变更：批量处理时显示改了哪些文件，先核对再保留",
                 "预览：在右侧面板内查看，无需下载即可判断好坏"],
     "note": "交付前一定过一遍结果区四视图。"},
    {"type": "content", "chapter": "第三章 核心功能逐个看", "kicker": "管理",
     "title": "侧边栏：多任务并行与管理",
     "bullets": ["任务列表按文件夹分组，支持搜索",
                 "可同时发起多个任务并行推进（如一边生成 PPT 一边做分析）",
                 "点开任一任务可继续追问、补充材料",
                 "底部显示当前账号头像与状态"],
     "note": "多任务并行是效率关键，鼓励同时跑独立子任务。"},

    # ═══════════ 第四章 进阶技巧 ═══════════
    {"type": "divider", "kicker": "第四章", "title": "进阶技巧",
     "note": "功能会用之后，再学两个决定成败的习惯。"},
    {"type": "content", "chapter": "第四章 进阶技巧", "kicker": "技巧 ①",
     "title": "先给方案，再执行",
     "bullets": ["目标不清就直接生成，结果必然跑偏",
                 "用 Plan 模式或提示词「先给方案 / 先出大纲」，确认后再动手",
                 "方案应包含：任务顺序、每步输入与输出、确认节点、交付物、风险",
                 "好处：把你的隐性要求显性化，减少返工"],
     "note": "这是新手最该养成的第一个习惯。"},
    {"type": "compare", "chapter": "第四章 进阶技巧", "kicker": "技巧 ①",
     "title": "对比：直接生成 vs 先给方案",
     "left_label": "❌ 直接生成", "left_items": [
         "「写一份 7 月复盘 PPT」",
         "WB 凭猜测定结构、选数据",
         "拿到后发现重点不对、数据口径错",
         "返工成本高，且不知改哪"],
     "right_label": "✅ 先给方案", "right_items": [
         "「先给 12 页逐页大纲，确认后再生成」",
         "你在校验点把关结构与口径",
         "终稿贴合需求，仅需局部微调",
         "一次到位，沟通成本最低"],
     "note": "对应提示词 01（先提问）与 07（综合 PPT 先大纲）。"},
    {"type": "flow", "chapter": "第四章 进阶技巧", "kicker": "技巧 ②",
     "title": "出错后怎么办：五步纠正工作流",
     "steps": ["WB 原始\n输出", "四查自检\n(数字/日期/\n人名/结论)", "定位错误\n(行号/出处)", "纠正话术\n(精确指出)", "人工终稿\n(通读导出)"],
     "note": "任何功能出错都套这五步，不要重头来。"},
    {"type": "compare", "chapter": "第四章 进阶技巧", "kicker": "技巧 ②",
     "title": "纠错示例：「待确认」被写成「决定」",
     "left_label": "❌ 错误", "left_items": [
         "原文「快闪预算追加金额待定」",
         "纪要却写成「决定追加预算 XX 元」",
         "AI 把模糊补成了确定结论",
         "后果：未经授权的承诺"],
     "right_label": "✅ 正确", "right_items": [
         "归入「待确认事项」并引用原话",
         "备注：缺具体金额，需后续确认",
         "原则：源材料未明确的标待确认",
         "后果：责任清晰、可追查"],
     "note": "这是文档生成功能最常见的幻觉类型。"},
    {"type": "content", "chapter": "第四章 进阶技巧", "kicker": "技巧 ③",
     "title": "写好提示词的四要素",
     "bullets": ["目标：要生成什么（纪要 / 报告 / PPT）",
                 "输入：分析哪些数据、参考哪些文档",
                 "输出：格式（Word/Excel/PPT）、页数、风格",
                 "约束：只据材料、不虚构、口径、时限、禁用项"],
     "note": "四要素越齐，WorkBuddy 越不需要猜。"},
    {"type": "content", "chapter": "第四章 进阶技巧", "kicker": "技巧 ④",
     "title": "数据一致性自检（四查清单）",
     "bullets": ["① 数字：每个值能在源表找到，合计 = 各项相加",
                 "② 日期：年月对、截止是未来、逻辑不自相矛盾",
                 "③ 人名：与参会表一致，负责人没搞反",
                 "④ 结论：有出处，「建议」≠「决定」，「待确认」≠ 确定"],
     "note": "详见 06-prompts/08-review-checklist.md 与 11-quality-assurance/。"},

    # ═══════════ 第五章 Demo 结果顺带展示 ═══════════
    {"type": "divider", "kicker": "第五章", "title": "Demo 结果顺带展示",
     "note": "前面已用流程讲清功能，这里只快速过一遍成果长什么样。"},
    {"type": "content", "chapter": "第五章 Demo 结果顺带展示", "kicker": "成果",
     "title": "青屿茶研案例 · 五类产出（一览）",
     "bullets": ["会议纪要：6 段结构（信息/讨论/决定/行动/待确认/下次）",
                 "工作总结：区分完成 / 参与 / 协助，延期项单列",
                 "活动通知：只含已确认事项",
                 "分析报告：事实 / 判断 / 建议分明，整体转化率 37.0%",
                 "复盘 PPT：12 页，每页一个结论"],
     "note": "这里不展开，成果在 05-demo-outputs/ 自行查看。"},
    {"type": "content", "chapter": "第五章 Demo 结果顺带展示", "kicker": "复现",
     "title": "一条命令复现全部 Demo",
     "bullets": ["所有 Office 文件由 tools/ 下 Python 脚本生成",
                 "数据来自单一源 case_data.py，含 verify() 自检",
                 "图表由 gen_charts.py 用 Pillow 绘制（含纠错标注）",
                 "仓库已发布到 GitHub，克隆即可复现"],
     "note": "把「功能演示」和「可复现」收口，呼应连接器功能⑧。"},

    # ═══════════ 第六章 收尾 ═══════════
    {"type": "divider", "kicker": "第六章", "title": "检查、复现与寄语",
     "note": "收尾，落到行动。"},
    {"type": "content", "chapter": "第六章 检查、复现与寄语", "kicker": "动作",
     "title": "三个核心动作",
     "bullets": ["先提问，再生成方案（功能：先提问 / Plan 模式）",
                 "上传材料，明确限制（功能：添加上下文 / 工作空间）",
                 "人工检查，再交付（功能：结果区验收 / 四查）"],
     "note": "让学员记住这三句话，对应前面学过的入口。"},
    {"type": "content", "chapter": "第六章 检查、复现与寄语", "kicker": "复现",
     "title": "GitHub 与二次开发",
     "bullets": ["仓库：github.com/onehount/workbuddy-beginner-training",
                 "README 给出三条复现路径（按角色选）",
                 "下载 inputs，按 06-prompts 自行跑各功能",
                 "二次开发请遵守虚构声明与双许可证"],
     "note": "演示仓库结构即可，呼应连接器功能⑧。"},
    {"type": "closing", "kicker": "寄语", "title": "写给刚起步的你",
     "bullets": ["初入职场，不必一次写出完美指令",
                 "理解工作本质，知道目标与判断标准",
                 "保持好奇，持续学习，与 AI 协同成长"],
     "note": "放慢语速，给学员留出思考空间。"},
]


# ============================================================
# PPTX 渲染
# ============================================================
def render_pptx(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    def footer(slide, idx, chapter):
        box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9), Inches(0.35))
        tf = box.text_frame; p = tf.paragraphs[0]; p.text = chapter or ""
        p.font.size = PptPt(11); p.font.color.rgb = GREY; p.font.name = "微软雅黑"
        box2 = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
        tf2 = box2.text_frame; p2 = tf2.paragraphs[0]; p2.text = f"{idx}/{total}"
        p2.font.size = PptPt(11); p2.font.color.rgb = GREY; p2.font.name = "微软雅黑"
        p2.alignment = PP_ALIGN.RIGHT

    def bg(slide, color):
        box = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        box.fill.solid(); box.fill.fore_color.rgb = color
        box.line.fill.background(); box.shadow.inherit = False
        return box

    def top_bar(slide, title, kicker):
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]; p.text = title
        p.font.size = PptPt(30); p.font.bold = True; p.font.color.rgb = WHITE
        p.font.name = "微软雅黑"
        sp = tf.add_paragraph(); sp.text = kicker
        sp.font.size = PptPt(13); sp.font.color.rgb = LIGHT; sp.font.name = "微软雅黑"

    def add_text(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.alignment = align
        p.font.size = PptPt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = "微软雅黑"
        return box

    def add_bullets(slide, l, t, w, h, items, size=20, color=GREY):
        if isinstance(items, str):
            items = [items]
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = "• " + it
            p.font.size = PptPt(size); p.font.color.rgb = color; p.font.name = "微软雅黑"
            p.space_after = PptPt(8)
        return box

    def rounded(slide, l, t, w, h, fill, line, lw=1.5):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = line; shp.line.width = PptPt(lw)
        shp.shadow.inherit = False
        return shp

    # —— 功能页：四象限 ——
    def add_feature_page(slide, s):
        bg(slide, WHITE)
        top_bar(slide, s["title"], s["kicker"])
        # 象限坐标
        cols = [(0.6, 5.95), (6.78, 5.95)]
        rows = [(1.45, 2.45), (4.05, 2.45)]
        specs = [("是什么", s["what"], BLUE),
                 ("干什么", s["do"], CYAN),
                 ("怎么用", s["how"], AMBER),
                 ("在哪里", s["where"], GREEN)]
        for ci in range(2):
            cx, cw = cols[ci]
            for rj in range(2):
                ry, ch = rows[rj]
                title, body, accent = specs[rj * 2 + ci]
                rounded(slide, cx, ry, cw, ch, LIGHT, BLUE, 1.5)
                # 标题条
                add_text(slide, cx + 0.2, ry + 0.12, cw - 0.4, 0.4,
                         title, 17, accent, bold=True)
                # 分隔线
                line = slide.shapes.add_shape(1, Inches(cx + 0.2), Inches(ry + 0.55),
                                              Inches(cw - 0.4), Inches(0.02))
                line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
                if isinstance(body, list):
                    add_bullets(slide, cx + 0.2, ry + 0.68, cw - 0.4, ch - 0.8,
                               body, size=15, color=INK)
                else:
                    add_text(slide, cx + 0.2, ry + 0.68, cw - 0.4, ch - 0.8,
                             body, 15, INK)
        # 底部 demo 提示条
        if s.get("demo"):
            rounded(slide, 0.6, 6.62, 12.13, 0.32, NAVY, NAVY, 1)
            add_text(slide, 0.75, 6.64, 11.8, 0.3, "🎯 " + s["demo"], 13, WHITE, bold=True)

    # —— 界面示意页 ——
    def add_mockup(slide, s):
        bg(slide, WHITE)
        top_bar(slide, s["title"], s["kicker"])
        # 侧边栏
        rounded(slide, 0.55, 1.5, 2.7, 5.2, NAVY, NAVY, 1)
        add_text(slide, 0.7, 1.62, 2.4, 0.4, "侧边栏", 16, CYAN, bold=True)
        add_bullets(slide, 0.7, 2.1, 2.45, 4.4,
                    ["任务列表（按文件夹）", "搜索任务", "选择工作空间", "当前账号头像"],
                    size=14, color=WHITE)
        # 对话区
        rounded(slide, 3.45, 1.5, 5.9, 5.2, LIGHT, BLUE, 1.5)
        add_text(slide, 3.62, 1.62, 5.6, 0.4, "对话区（核心交互）", 16, NAVY, bold=True)
        add_bullets(slide, 3.62, 2.1, 5.6, 2.6,
                    ["任务标题栏", "消息列表（执行过程）", "输入框：一句话下达任务",
                     "@ 引用 / 粘贴截图 / 上传 / 补充", "切换工作模式 Craft/Plan/Ask"],
                    size=14, color=INK)
        # 对话区里的输入框示意
        inp = rounded(slide, 3.62, 5.5, 5.55, 0.95, WHITE, BLUE, 1.5)
        add_text(slide, 3.75, 5.62, 5.3, 0.7,
                 "「把这份会议记录整理成纪要，先给结构」", 13, GREY)
        # 结果区
        rounded(slide, 9.55, 1.5, 3.25, 5.2, PptRGB(0xE8, 0xF8, 0xEE), GREEN, 1.5)
        add_text(slide, 9.72, 1.62, 3.0, 0.4, "结果区（右侧）", 16, GREEN, bold=True)
        add_bullets(slide, 9.72, 2.1, 3.0, 4.4,
                    ["产物：预览生成文件", "全部文件：产物清单", "变更：批量改了什么",
                     "预览：面板内查看"],
                    size=14, color=INK)
        # 底注
        add_text(slide, 0.6, 6.85, 12, 0.4,
                 "在哪里：左侧栏管任务，中间对话区下达与追问，右侧结果区验收产物。",
                 14, BLUE, bold=True)

    for idx, s in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        t = s["type"]

        if t == "cover":
            bg(slide, NAVY)
            add_text(slide, 1.0, 2.0, 11.3, 0.6, s["kicker"], 22, LIGHT)
            add_text(slide, 1.0, 2.6, 11.3, 1.2, s["title"], 42, WHITE, bold=True)
            add_text(slide, 1.0, 4.1, 11.3, 2.0,
                     "\n".join("· " + b for b in s["bullets"]), 19, LIGHT)
            footer(slide, idx, "WorkBuddy 零基础办公任务实战教程")

        elif t == "divider":
            bg(slide, NAVY)
            add_text(slide, 1.0, 2.6, 11.3, 0.8, s["kicker"], 26, CYAN, bold=True)
            add_text(slide, 1.0, 3.4, 11.3, 1.4, s["title"], 46, WHITE, bold=True)
            footer(slide, idx, s["title"])

        elif t == "content":
            bg(slide, WHITE)
            top_bar(slide, s["title"], s["kicker"])
            add_bullets(slide, 0.8, 1.6, 11.7, 5.2, s["bullets"], size=21)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "feature":
            add_feature_page(slide, s)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "mockup":
            add_mockup(slide, s)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "chart":
            bg(slide, WHITE)
            top_bar(slide, s["title"], s["kicker"])
            img_path = os.path.join(CHART_DIR, s["image"])
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.45), width=Inches(9.2))
                add_text(slide, 0.8, 6.0, 9.0, 0.5, s.get("caption", ""), 14, BLUE)
            else:
                add_text(slide, 0.8, 2.0, 10, 3, f"[图表缺失: {s['image']}]", 18, RED)
            if s.get("bullets"):
                add_bullets(slide, 10.0, 1.6, 3.0, 4.8, s["bullets"], size=16)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "compare":
            bg(slide, WHITE)
            top_bar(slide, s["title"], s["kicker"])
            lh = min(len(s["left_items"]) * 0.62 + 0.8, 4.9)
            left_box = rounded(slide, 0.6, 1.5, 5.9, lh, PptRGB(0xFD, 0xED, 0xEC), RED, 1.5)
            add_text(slide, 0.85, 1.58, 5.5, 0.45, s["left_label"], 18, RED, bold=True)
            add_bullets(slide, 0.85, 2.1, 5.5, lh - 0.5, s["left_items"], size=16, color=RED)
            rh = min(len(s["right_items"]) * 0.62 + 0.8, 4.9)
            rounded(slide, 6.8, 1.5, 5.9, rh, PptRGB(0xE8, 0xF8, 0xEE), GREEN, 1.5)
            add_text(slide, 7.05, 1.58, 5.5, 0.45, s["right_label"], 18, GREEN, bold=True)
            add_bullets(slide, 7.05, 2.1, 5.5, rh - 0.5, s["right_items"], size=16, color=GREEN)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "flow":
            bg(slide, WHITE)
            top_bar(slide, s["title"], s["kicker"])
            steps = s["steps"]; n = len(steps)
            sw = 1.95; gap = 0.32
            total_w = n * sw + (n - 1) * gap
            sx = (13.333 - total_w) / 2; sy = 2.6; sh = 2.6
            for i, label in enumerate(steps):
                x = sx + i * (sw + gap)
                box = rounded(slide, x, sy, sw, sh,
                              NAVY if i == 0 else (GREEN if i == n - 1 else LIGHT), BLUE, 1.5)
                txt_color = WHITE if i in (0, n - 1) else INK
                tb = slide.shapes.add_textbox(Inches(x + 0.08), Inches(sy + 0.15),
                                              Inches(sw - 0.16), Inches(sh - 0.3))
                tf = tb.text_frame; tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]; p.text = label
                p.font.size = PptPt(14); p.font.color.rgb = txt_color
                p.font.bold = True; p.font.name = "微软雅黑"; p.alignment = PP_ALIGN.CENTER
                if i < n - 1:
                    ax = x + sw + 0.04
                    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                                 Inches(ax), Inches(sy + sh / 2 - 0.18),
                                                 Inches(gap - 0.08), Inches(0.36))
                    arr.fill.solid(); arr.fill.fore_color.rgb = CYAN; arr.line.fill.background()
            add_text(slide, 0.8, 5.6, 11.7, 0.9,
                     "每步对应界面操作：新建任务 → 上传/输入 → 写约束 → 确认大纲 → 结果区导出。",
                     14, BLUE)
            footer(slide, idx, s.get("chapter", ""))

        elif t == "closing":
            bg(slide, NAVY)
            add_text(slide, 1.0, 2.0, 11.3, 0.6, s["kicker"], 22, CYAN, bold=True)
            add_text(slide, 1.0, 2.7, 11.3, 1.0, s["title"], 40, WHITE, bold=True)
            add_text(slide, 1.0, 3.9, 11.3, 2.4,
                     "\n".join("· " + b for b in s["bullets"]), 22, LIGHT)
            footer(slide, idx, "寄语")

    prs.save(path)
    return path


# ============================================================
# PDF 渲染（reportlab，与 PPT 内容一致）
# ============================================================
def render_pdf(path):
    pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
    FONT = 'STSong-Light'
    pw, ph = A4[0], A4[1]
    c = RLCanvas.Canvas(path, pagesize=(pw, ph))
    total = len(SLIDES)

    def rect(x, y, w, h, fill, stroke=None):
        if fill:
            c.setFillColorRGB(*[v / 255 for v in fill]); c.rect(x, y, w, h, fill=1, stroke=0)
        if stroke:
            c.setStrokeColorRGB(*[v / 255 for v in stroke]); c.setLineWidth(2)
            c.rect(x, y, w, h, fill=0, stroke=1)

    def text(x, y, s, size, color, bold=False):
        c.setFillColorRGB(*[v / 255 for v in color]); c.setFont(FONT, size)
        c.drawString(x, y, s)

    def footer_p(idx, chapter):
        c.setFillColorRGB(*[v / 255 for v in GREY_H]); c.setFont(FONT, 10)
        c.drawString(30, 22, chapter); c.drawRightString(pw - 30, 22, f"{idx}/{total}")

    def bullets(items, x, y, size, color, lh=26, max_chars=34):
        if isinstance(items, str):
            items = [items]
        yy = y
        for it in items:
            while len(it) > max_chars:
                c.drawString(x, yy, "• " + it[:max_chars] + "…"); yy -= lh
                it = "…" + it[max_chars:]
            c.drawString(x, yy, "• " + it); yy -= lh
        return yy

    for idx, s in enumerate(SLIDES, start=1):
        t = s["type"]
        if t == "cover":
            rect(0, 0, pw, ph, NAVY_H)
            text(60, ph - 150, s["kicker"], 20, (0xEA, 0xF1, 0xFD))
            text(60, ph - 220, s["title"], 36, (0xFF, 0xFF, 0xFF), bold=True)
            yy = ph - 300
            for b in s["bullets"]:
                text(70, yy, "· " + b, 17, (0xEA, 0xF1, 0xFD)); yy -= 30
            footer_p(idx, "WorkBuddy 零基础办公任务实战教程")

        elif t == "divider":
            rect(0, 0, pw, ph, NAVY_H)
            text(60, ph - 220, s["kicker"], 24, (0x17, 0xC0, 0xC7), bold=True)
            text(60, ph - 300, s["title"], 42, (0xFF, 0xFF, 0xFF), bold=True)
            footer_p(idx, s["title"])

        elif t == "content":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 24, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            bullets(s["bullets"], 60, ph - 130, 18, GREY_H)
            footer_p(idx, s.get("chapter", ""))

        elif t == "feature":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 23, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            # 四象限
            quads = [("是什么", s["what"], BLUE_H), ("干什么", s["do"], (0x17, 0xC0, 0xC7)),
                     ("怎么用", s["how"], (0xE6, 0x7E, 0x22)), ("在哪里", s["where"], GREEN_H)]
            positions = [(40, ph - 430, pw / 2 - 60, 150),
                         (pw / 2 + 20, ph - 430, pw / 2 - 60, 150),
                         (40, ph - 600, pw / 2 - 60, 150),
                         (pw / 2 + 20, ph - 600, pw / 2 - 60, 150)]
            for (qt, qb, qc), (qx, qy, qw, qh) in zip(quads, positions):
                rect(qx, qy, qw, qh, (0xEA, 0xF1, 0xFD), stroke=BLUE_H)
                text(qx + 12, qy + qh - 24, qt, 15, qc, bold=True)
                if isinstance(qb, list):
                    bullets(qb, qx + 12, qy + qh - 48, 13, GREY_H, lh=20, max_chars=22)
                else:
                    text(qx + 12, qy + qh - 48, qb, 13, GREY_H)
            if s.get("demo"):
                text(40, ph - 640, "🎯 " + s["demo"], 13, GREEN_H, bold=True)
            footer_p(idx, s.get("chapter", ""))

        elif t == "mockup":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 23, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            text(50, ph - 130, "左侧边栏：任务列表 / 搜索 / 选择工作空间 / 账号头像", 15, NAVY_H, bold=True)
            text(50, ph - 160, "中对话区：任务标题栏 / 消息列表 / 输入框(一句话下达任务) / @引用·粘贴·上传 / 模式切换", 14, GREY_H)
            text(50, ph - 190, "右结果区：产物 / 全部文件 / 变更 / 预览", 15, GREEN_H, bold=True)
            text(50, ph - 230, "在哪里：左侧栏管任务，中间下达与追问，右侧验收产物。", 14, BLUE_H, bold=True)
            footer_p(idx, s.get("chapter", ""))

        elif t == "chart":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 22, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            img_path = os.path.join(CHART_DIR, s["image"])
            if os.path.exists(img_path):
                from reportlab.lib.utils import ImageReader
                ir = ImageReader(img_path); iw, ih = ir.getSize()
                scale = min((pw - 100) / iw, (ph - 250) / ih) * 0.88
                c.drawImage(ir, 50, ph - 430, iw * scale, ih * scale)
            if s.get("caption"):
                text(50, ph - 450, s["caption"], 13, BLUE_H)
            if s.get("bullets"):
                bullets(s.get("bullets"), 540, ph - 130, 14, GREY_H, lh=24)
            footer_p(idx, s.get("chapter", ""))

        elif t == "compare":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 22, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            rect(40, ph - 420, pw / 2 - 60, 310, (0xFD, 0xED, 0xEC), stroke=RED_H)
            text(52, ph - 400, s["left_label"], 16, RED_H, bold=True)
            bullets(s["left_items"], 56, ph - 430, 14, RED_H, lh=22)
            rect(pw / 2 + 20, ph - 420, pw / 2 - 60, 310, (0xE8, 0xF8, 0xEE), stroke=GREEN_H)
            text(pw / 2 + 32, ph - 400, s["right_label"], 16, GREEN_H, bold=True)
            bullets(s["right_items"], pw / 2 + 36, ph - 430, 14, GREEN_H, lh=22)
            footer_p(idx, s.get("chapter", ""))

        elif t == "flow":
            rect(0, 0, pw, ph, (0xFF, 0xFF, 0xFF))
            rect(0, ph - 90, pw, 90, NAVY_H)
            text(50, ph - 62, s["title"], 22, (0xFF, 0xFF, 0xFF), bold=True)
            text(50, ph - 82, s["kicker"], 12, (0xEA, 0xF1, 0xFD))
            steps = s["steps"]; n = len(steps)
            bw = (pw - 120) / n - 20; bh = 150; bx0 = 60; by = ph - 320
            for i, lbl in enumerate(steps):
                bx = bx0 + i * (bw + 20)
                fc = NAVY_H if i == 0 else (GREEN_H if i == n - 1 else LIGHT_H)
                tc = (0xFF, 0xFF, 0xFF) if i in (0, n - 1) else INK
                rect(bx, by, bw, bh, fc, stroke=BLUE_H)
                lines = lbl.split("\n"); cy = by + bh - 25
                for ln in lines:
                    c.setFillColorRGB(*[v / 255 for v in tc])
                    c.setFont(FONT, 12); c.drawCentredString(bx + bw / 2, cy, ln); cy -= 18
                if i < n - 1:
                    c.setFillColorRGB(*[v / 255 for v in (0x17, 0xC0, 0xC7)])
                    c.drawString(bx + bw + 3, by + bh / 2 - 5, "→")
            text(50, ph - 450, "每步对应界面操作：新建任务 → 上传/输入 → 写约束 → 确认大纲 → 结果区导出。", 13, BLUE_H)
            footer_p(idx, s.get("chapter", ""))

        elif t == "closing":
            rect(0, 0, pw, ph, NAVY_H)
            text(60, ph - 160, s["kicker"], 22, (0x17, 0xC0, 0xC7), bold=True)
            text(60, ph - 220, s["title"], 34, (0xFF, 0xFF, 0xFF), bold=True)
            yy = ph - 290
            for b in s["bullets"]:
                text(70, yy, "· " + b, 19, (0xEA, 0xF1, 0xFD)); yy -= 30
            footer_p(idx, "寄语")

        c.showPage()
    c.save()
    return path


# ============================================================
# 大纲 / 备注
# ============================================================
def write_outline(path):
    lines = ["# 课程页面索引（slide-outline.md）", "",
             "> 与 `workbuddy-beginner-training.pptx` / `.pdf` 严格对应。", "",
             "> 本版本为「功能导向」：先讲各功能是什么/干什么/怎么用/在哪里，",
             "> 用统一案例演示使用流程，再加进阶技巧，Demo 结果顺带展示。", ""]
    type_cn = {"cover": "封面", "divider": "章节", "content": "内容", "feature": "功能",
               "mockup": "界面示意", "chart": "图表", "compare": "对比", "flow": "流程", "closing": "结语"}
    for i, s in enumerate(SLIDES, start=1):
        extra = ""
        if s.get("image"):
            extra += f" | 图:{s['image']}"
        if s.get("steps"):
            extra += f" | 流程({len(s['steps'])}步)"
        if s["type"] == "feature":
            extra += " | 功能四象限"
        lines.append(f"{i:02d} | {type_cn[s['type']]} | {s['title']}{extra}")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")
    return path


def write_notes(path):
    lines = ["# 讲师备注（slide-notes.md）", "",
             "> 每页一句到两句的讲解要点，配合 `full-speaking-script.md` 使用。", ""]
    for i, s in enumerate(SLIDES, start=1):
        lines.append(f"## 第 {i} 页 · {s['title']}")
        lines.append("")
        lines.append(s.get("note", ""))
        lines.append("")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")
    return path


if __name__ == "__main__":
    os.makedirs(SLIDE_DIR, exist_ok=True)
    p1 = render_pptx(os.path.join(SLIDE_DIR, "workbuddy-beginner-training.pptx"))
    p2 = render_pdf(os.path.join(SLIDE_DIR, "workbuddy-beginner-training.pdf"))
    p3 = write_outline(os.path.join(SLIDE_DIR, "slide-outline.md"))
    p4 = write_notes(os.path.join(SLIDE_DIR, "slide-notes.md"))
    print(f"[OK] {os.path.relpath(p1, ROOT)}  ({len(SLIDES)} 页)")
    print(f"[OK] {os.path.relpath(p2, ROOT)}")
    print(f"[OK] {os.path.relpath(p3, ROOT)}")
    print(f"[OK] {os.path.relpath(p4, ROOT)}")
