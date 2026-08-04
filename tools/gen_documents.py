# -*- coding: utf-8 -*-
"""
生成 Demo 输入文档与模板（04-demo-inputs/）

- documents/   : 01_案例背景说明 / 02_会议原始记录 / 06_个人工作记录
- templates/   : 07_会议纪要模板 / 08_工作总结模板 / 09_通知模板
                 / 10_分析报告模板 / 11_月度复盘PPT模板

所有数字来自 case_data.py；作者统一为 "WorkBuddy Training Project"；
每页均带虚构声明。

License: MIT
"""
import os
import sys

from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import case_data as C

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
DOC_DIR = os.path.join(ROOT, "04-demo-inputs", "documents")
TPL_DIR = os.path.join(ROOT, "04-demo-inputs", "templates")

NAVY = RGBColor(0x1B, 0x3A, 0x6B)
BLUE = RGBColor(0x2E, 0x6B, 0xE6)
GREY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xC0, 0x39, 0x2B)
FONT = "微软雅黑"


# ============================================================
# 通用文档工具
# ============================================================
def set_run(run, size=11, bold=False, color=None, name=FONT):
    run.font.name = name
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn('w:rFonts'))
    if rfonts is None:
        rfonts = OxmlElement('w:rFonts')
        rpr.append(rfonts)
    rfonts.set(qn('w:eastAsia'), name)
    rfonts.set(qn('w:ascii'), name)
    rfonts.set(qn('w:hAnsi'), name)
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def set_style_font(doc):
    style = doc.styles['Normal']
    style.font.name = FONT
    style.font.size = Pt(11)
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn('w:rFonts'))
    if rfonts is None:
        rfonts = OxmlElement('w:rFonts')
        rpr.append(rfonts)
    rfonts.set(qn('w:eastAsia'), FONT)
    rfonts.set(qn('w:ascii'), FONT)
    rfonts.set(qn('w:hAnsi'), FONT)


def add_title(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    set_run(r, size=18, bold=True, color=NAVY)
    p.space_after = Pt(6)
    return p


def add_h1(doc, text):
    p = doc.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run(text)
    set_run(r, size=14, bold=True, color=NAVY)
    return p


def add_h2(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    set_run(r, size=12, bold=True, color=BLUE)
    return p


def add_body(doc, text, size=11):
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = 1.25
    r = p.add_run(text)
    set_run(r, size=size)
    return p


def add_note(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    set_run(r, size=9, color=GREY)
    return p


def add_bullet(doc, text, size=11):
    p = doc.add_paragraph(style='List Bullet')
    r = p.add_run(text)
    set_run(r, size=size)
    return p


def style_table_header(row):
    for cell in row.cells:
        for p in cell.paragraphs:
            for r in p.runs:
                set_run(r, size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        tcPr = cell._tc.get_or_add_tcPr()
        shd = tcPr.find(qn('w:shd'))
        if shd is None:
            shd = OxmlElement('w:shd')
            tcPr.append(shd)
        shd.set(qn('w:fill'), "2E6BE6")
        shd.set(qn('w:val'), "clear")
        shd.set(qn('w:color'), "auto")


def set_cell(cell, text, size=10, bold=False, color=None):
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(text)
    set_run(r, size=size, bold=bold, color=color)


def set_core_props(doc, title):
    cp = doc.core_properties
    cp.author = C.AUTHOR
    cp.last_modified_by = C.AUTHOR
    cp.title = title
    cp.subject = "WorkBuddy 零基础教学项目 · 虚构演示材料"
    cp.comments = C.DISCLAIMER
    cp.category = "Training Material"
    cp.revision = 1


# ============================================================
# 01_案例背景说明.docx
# ============================================================
def gen_background():
    doc = Document()
    set_style_font(doc)
    add_title(doc, f"{C.BRAND_FULL} · 市场部 {C.PERIOD_LABEL} 业务复盘案例背景")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")
    add_note(doc, f"本文档用于 WorkBuddy 零基础教学演示，不对应任何真实企业。")

    add_h1(doc, "一、品牌简介")
    add_body(doc,
        f"{C.BRAND_CN}（{C.BRAND_EN}）是一个虚构的新式茶饮品牌，主打原叶茶与季节限定饮品。"
        "本教学案例以该品牌市场部的一次月度业务复盘为场景，演示如何用 WorkBuddy 把零散材料"
        "整理成正式办公成果。")

    add_h1(doc, "二、部门与人员角色")
    add_body(doc, f"本次复盘由{C.BRAND_CN}{C.DEPARTMENT}组织，共 {len(C.PEOPLE)} 人参会：")
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = 'Table Grid'
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(["姓名", "岗位", "所属", "本月主要负责"]):
        set_cell(tbl.rows[0].cells[i], h)
    style_table_header(tbl.rows[0])
    for (name, role, dept, duty) in C.PEOPLE:
        cells = tbl.add_row().cells
        set_cell(cells[0], name.replace("　", ""), bold=True)
        set_cell(cells[1], role)
        set_cell(cells[2], dept)
        set_cell(cells[3], duty)

    add_h1(doc, "三、本月市场活动")
    add_body(doc, f"{C.PERIOD_LABEL}共开展 {C.CAMPAIGN_TOTAL} 项市场活动，其中"
        f"{C.CAMPAIGN_DONE} 项已完成，{C.CAMPAIGN_DELAYED} 项延期。")
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["活动名称", "开始", "计划结束", "状态", "负责人"]):
        set_cell(tbl.rows[0].cells[i], h)
    style_table_header(tbl.rows[0])
    for (idx, name, start, end, plan_end, status, owner, helper, ch, desc) in C.CAMPAIGNS:
        cells = tbl.add_row().cells
        set_cell(cells[0], name)
        set_cell(cells[1], start)
        set_cell(cells[2], plan_end)
        set_cell(cells[3], status, bold=True,
                 color=(RED if status == "延期" else NAVY))
        set_cell(cells[4], owner.replace("　", ""))

    add_h1(doc, "四、渠道与数据概览")
    add_body(doc,
        f"本月从 {len(C.CHANNELS)} 个渠道获取曝光与线索。统一计算口径：{C.METRIC_FORMULA}。"
        "整体转化率 = 有效线索合计 ÷ 线索合计，不取各渠道转化率的算术平均值。")
    t = C.channel_totals()
    e, l, q, r = C.grand_total()
    add_bullet(doc, f"总曝光量：{e:,} 次")
    add_bullet(doc, f"总线索数量：{l:,} 条")
    add_bullet(doc, f"总有效线索：{q:,} 条")
    add_bullet(doc, f"整体转化率：{r:.1f}%")
    add_note(doc, "明细数据见 05_渠道曝光与线索数据.xlsx（汇总值均由公式得出）。")

    add_h1(doc, "五、复盘会安排")
    add_bullet(doc, f"会议名称：{C.MEETING_TITLE}")
    add_bullet(doc, f"时间：{C.MEETING_DATE} {C.MEETING_TIME}")
    add_bullet(doc, f"地点：{C.MEETING_PLACE}")

    path = os.path.join(DOC_DIR, "01_案例背景说明.docx")
    set_core_props(doc, "案例背景说明")
    doc.save(path)
    return path


# ============================================================
# 02_市场部月度会议原始记录.docx
# ============================================================
MEETING_RAW = """青屿茶研市场部 2026 年 7 月月度业务复盘会 —— 原始记录（未经整理）

会议时间：2026年8月3日 14:00–15:40
会议地点：青屿茶研总部 3 楼竹里会议室
主持人：沈知远（市场部负责人）
参会：沈知远、陈屿、李维、周然、苏眠、郑楠
记录方式：现场速记，未署具体记录人

[14:00 会议开始]
沈知远：好，人都到齐了，我们开始吧。今天这个会主要是把七月份的市场活动过一遍，看看做得怎么样，有哪些东西要接着干。七月份我们一共搞了五个活动，对吧？
陈屿：对，五个。夏日新品试饮、小红书种草、微信社群优惠、商场快闪，还有一个校园联名。
沈知远：嗯，那个校园联名我看状态有点问题，到底做没做完？
苏眠：这个……校园联名它没按期做完。本来计划七月底结束的，但是合作方那个 Logo 使用规范一直没给我们确认，我们也不敢随便把物料定稿，所以就延了。现在算延期。
沈知远：行，那这个就不算七月份完成的。纪要里要写清楚，校园联名是延期，不计入七月份完成率考核。这个大家没意见吧？
（众人点头）
陈屿：没意见，这个本来就没做完。

[14:10 市场活动情况]
沈知远：那先过一下前面四个做完了的。试饮那个，苏眠你讲讲。
苏眠：夏日新品试饮是七月五号到十二号，覆盖十二家直营门店，主推两款新品，一个青梅乌龙，一个桂香冷萃，现场试饮。反馈还不错，回头我把偏好比例整理一下给大家。
陈屿：快闪那个我补充一下，七月十八到二十号，三个核心商圈，我们还同步做了抖音本地推流。效果……反正现场人气是可以的，客流量还行。
周然：小红书种草是八号到二十五号，我们找了十八个达人，出了二十四篇笔记。李维那边投的，达人名单是苏眠之前帮我对了客群画像的。
李维：对，主要是小红书和抖音两个号在运营，日常内容也我这边发。抖音七月曝光还挺大的。

[14:25 渠道数据]
沈知远：好，那说数据。郑楠你把渠道数据拉一下。
郑楠：我这边统计的是五个渠道，小红书、抖音、微信公众号、微信社群，还有线下门店。七月整月总曝光五十八万三千，总线索四千七百一十，有效线索一千七百四十四。
李维：等一下，抖音曝光这么高，转化率是不是偏低？我总觉得抖音那边的量上去了，但质量好像一般。
郑楠：对，抖音曝光二十六万二，线索一千五百六，但有效线索只有三百九，转化率百分之二十五，是全渠道最低的。
沈知远：那小红书呢？
郑楠：小红书曝光十八万五，线索一千二百八，有效线索四百三十六，转化率三十四点一。微信社群最高，曝光三万二，但转化率有百分之六十。
陈屿：社群虽然量小，但质量高啊，都是门店沉淀的老客。
郑楠：所以整体转化率我们是按有效线索除以总线索算的，一千七百四十四除以四千七百一十，大概百分之三十七。
沈知远：等等，这个口径大家要统一一下。以前各个渠道自己算自己的，乱。从八月份开始，全部统一用"有效线索除以线索乘以百分之百"，各渠道不准再自己定义。这个要写进纪要，当作决定。
周然：同意，不然每次对数据都对不上，财务那边的口径也不一样。

[14:40 抖音与社群讨论]
李维：那抖音八月怎么办？预算还加不加？我这边想再投一点本地推。
沈知远：抖音八月先别动预算，你先做内容优化，把内容和目标人群匹配度再核一遍。预算暂时不调。
周然：那我跟你一起做 A/B 测试素材，我做两组。
李维：行，那我出匹配度复核报告，八月十二号前给你。
沈知远：好。社群那边，苏眠你八月搞一次扩容测试，验证一下转化率会不会随着社群规模下降。
苏眠：可以，我覆盖五个门店社群，八月十号前把方案拟出来。

[14:55 待确认事项]
沈知远：还有几个事没定下来。快闪那边，陈屿你不是说要追加点预算吗？
陈屿：对，快闪那边再加一点吧，我回头跟财务说。
沈知远：加多少？这个得有个数，而且要走财务审批，你定个金额。
陈屿：这个……我还没算好，回头弄，反正先记着要追加。
郑楠：还有那个数据看板，上次说要尽快弄出来，具体什么时候上线？
沈知远：看板尽快，但今天定不了时间，下周给节点吧。
李维：另外八月要不要上新短视频平台，我跟陈屿意见不太一样，今天先不拍板，下次再说。
陈屿：对，这个会上先放着，不急。
苏眠：校园联名第二阶段要不要扩展到第三所高校？如果第一批效果好可以再看看。
沈知远：这个也先不决定，看第一批情况。还有下半年品牌代言人那个方向，我们下次单独开个会聊，今天不展开了。

[15:20 收尾]
沈知远：行，那纪要苏眠或者周然整理一下，行动事项我念一遍大家确认：一，李维出抖音匹配度报告，八月十二；二，周然做 A/B 素材两组，八月十五；三，苏眠校园联名物料定稿八月八号；四，苏眠社群扩容方案八月十号；五，郑楠汇总七月数据归档，八月七号；六，陈屿起草八月活动通知，八月六号，要经我确认再发。大家都记一下。
陈屿：收到。
苏眠：好，我整理。
沈知远：好，那就这样，散会。
"""


def gen_meeting_raw():
    doc = Document()
    set_style_font(doc)
    add_title(doc, f"{C.MEETING_TITLE} · 原始记录")
    add_note(doc, "说明：本记录为现场速记稿，含口语重复、未决事项与待确认内容，未经整理。"
                  "正式纪要以 05-demo-outputs 中 WorkBuddy 生成并经人工审核的版本为准。")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")

    for line in MEETING_RAW.split("\n"):
        if not line.strip():
            continue
        if line.startswith("[") or line.endswith("]") or line.startswith("（") or line.startswith("("):
            p = doc.add_paragraph()
            r = p.add_run(line)
            set_run(r, size=10, color=GREY)
        else:
            p = doc.add_paragraph()
            p.paragraph_format.line_spacing = 1.3
            r = p.add_run(line)
            set_run(r, size=11)

    path = os.path.join(DOC_DIR, "02_市场部月度会议原始记录.docx")
    set_core_props(doc, "市场部月度会议原始记录")
    doc.save(path)
    return path


# ============================================================
# 06_市场专员个人工作记录.docx
# ============================================================
def gen_work_log():
    doc = Document()
    set_style_font(doc)
    add_title(doc, f"{C.BRAND_CN} {C.DEPARTMENT} · {C.PROTAGONIST} {C.PERIOD_LABEL} 个人工作记录")
    add_note(doc, f"记录人：{C.PROTAGONIST}（{C.PROTAGONIST_ROLE}）　统计周期：{C.PERIOD_RANGE}")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")

    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["序号", "日期", "工作内容", "类型"]):
        set_cell(tbl.rows[0].cells[i], h)
    style_table_header(tbl.rows[0])
    for i, (date, content, kind) in enumerate(C.WORK_LOG, start=1):
        cells = tbl.add_row().cells
        set_cell(cells[0], str(i), bold=True)
        set_cell(cells[1], date)
        set_cell(cells[2], content)
        set_cell(cells[3], kind,
                 color=(BLUE if kind == "协同" else NAVY), bold=True)

    add_note(doc, "「类型」说明：独立完成的事项标注「独立」，参与或协助团队完成的事项标注「协同」。"
                  "生成工作总结时，须据此区分个人成果与团队成果。")

    path = os.path.join(DOC_DIR, "06_市场专员个人工作记录.docx")
    set_core_props(doc, "市场专员个人工作记录")
    doc.save(path)
    return path


# ============================================================
# 07_会议纪要模板.docx
# ============================================================
def gen_tpl_minutes():
    doc = Document()
    set_style_font(doc)
    add_title(doc, "会议纪要模板")
    add_note(doc, "使用说明：将本模板与《市场部月度会议原始记录》一同上传给 WorkBuddy，"
                  "要求按原始材料整理，不得虚构。本模板为占位结构，填写处以「（请填写）」标出。")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")

    add_h2(doc, "会议基本信息")
    for label in ["会议名称", "会议时间", "会议地点", "主持人", "记录人", "参会人员"]:
        p = doc.add_paragraph()
        r = p.add_run(f"{label}：（请填写）"); set_run(r, size=11)

    add_h1(doc, "一、会议背景与目标")
    add_body(doc, "（请填写：会议起因、要解决的问题、期望达成的结论）")

    add_h1(doc, "二、主要讨论内容")
    add_body(doc, "（请填写：按议题归纳讨论要点，删除口语重复，保留真实含义）")

    add_h1(doc, "三、已确认决定")
    add_body(doc, "（请填写：明确形成的决定，逐条列出）")

    add_h1(doc, "四、行动事项")
    tbl = doc.add_table(rows=1, cols=6)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["序号", "事项", "负责人", "协同", "截止日期", "备注"]):
        set_cell(tbl.rows[0].cells[i], h)
    style_table_header(tbl.rows[0])
    for n in range(1, 5):
        cells = tbl.add_row().cells
        for i, _ in enumerate(["序号", "事项", "负责人", "协同", "截止日期", "备注"]):
            set_cell(cells[i], "（请填写）", size=10)

    add_h1(doc, "五、待确认事项")
    add_body(doc, "（请填写：原始记录中未明确负责人、金额、日期或结论的内容，逐条标记）")

    add_h1(doc, "六、下次会议安排")
    add_body(doc, "（请填写：时间、议题、负责人）")

    path = os.path.join(TPL_DIR, "07_会议纪要模板.docx")
    set_core_props(doc, "会议纪要模板")
    doc.save(path)
    return path


# ============================================================
# 08_月度工作总结模板.docx
# ============================================================
def gen_tpl_summary():
    doc = Document()
    set_style_font(doc)
    add_title(doc, "月度工作总结模板")
    add_note(doc, "使用说明：上传个人工作记录、会议纪要、活动完成情况，由 WorkBuddy 生成。"
                  "独立完成用「完成」，参与/协助用「参与」或「协助」，不得把团队成果写成个人成果。")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")

    add_h2(doc, "基本信息")
    for label in ["姓名", "岗位", "统计周期"]:
        p = doc.add_paragraph(); r = p.add_run(f"{label}：（请填写）"); set_run(r, size=11)

    for sec in ["一、本月已完成工作", "二、主要成果", "三、正在推进的工作",
                "四、存在的问题", "五、下月计划", "六、需要支持的事项"]:
        add_h1(doc, sec)
        add_body(doc, "（请填写）")

    path = os.path.join(TPL_DIR, "08_月度工作总结模板.docx")
    set_core_props(doc, "月度工作总结模板")
    doc.save(path)
    return path


# ============================================================
# 09_活动通知模板.docx
# ============================================================
def gen_tpl_notice():
    doc = Document()
    set_style_font(doc)
    add_title(doc, "活动通知模板")
    add_note(doc, "使用说明：依据已确认的会议行动事项生成通知。不得把待确认事项写成正式决定。"
                  "明确对象、时间、要求与反馈方式。")
    add_note(doc, f"⚠️ {C.DISCLAIMER}")

    p = doc.add_paragraph(); r = p.add_run("通知标题：（请填写）"); set_run(r, size=13, bold=True, color=NAVY)
    p = doc.add_paragraph(); r = p.add_run("通知对象：（请填写）"); set_run(r, size=11)

    add_h1(doc, "一、活动背景")
    add_body(doc, "（请填写：一句话说明背景，不重新解释全部会议过程）")

    add_h1(doc, "二、主要活动安排")
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for i, h in enumerate(["活动/任务", "负责人", "时间节点", "材料提交要求", "反馈方式"]):
        set_cell(tbl.rows[0].cells[i], h)
    style_table_header(tbl.rows[0])
    for _ in range(3):
        cells = tbl.add_row().cells
        for i in range(5):
            set_cell(cells[i], "（请填写）", size=10)

    add_h1(doc, "三、反馈与联系方式")
    add_body(doc, "（请填写：截止时间、联系人、提交渠道）")

    path = os.path.join(TPL_DIR, "09_活动通知模板.docx")
    set_core_props(doc, "活动通知模板")
    doc.save(path)
    return path


# ============================================================
# 10_市场分析报告模板.docx
# ============================================================
def gen_tpl_report():
    doc = Document()
    set_style_font(doc)
    add_title(doc, "市场分析报告模板")
    add_note(doc, "使用说明：分析渠道曝光与线索数据。必须区分「数据事实 / 分析判断 / 改进建议」。"
                  "所有比例说明计算口径；不虚构预算、销售额或用户画像。")
    add_note(doc, f"⚠️ {C.DISCLAIMER}　统一口径：{C.METRIC_FORMULA}")

    sections = [
        ("一、数据完整性检查", "（请填写：缺失项、异常值、口径一致性）"),
        ("二、市场活动完成情况", "（请填写：数量、完成率、延期说明）"),
        ("三、各渠道曝光与线索", "（请填写：表格或分渠道说明）"),
        ("四、转化率分析", "（请填写：各渠道转化率与整体转化率，注明公式）"),
        ("五、表现较好的渠道", "（请填写：依据数据事实）"),
        ("六、需要关注的问题", "（请填写：依据数据事实）"),
        ("七、数据可支持的改进建议", "（请填写：仅限数据支撑的建议）"),
    ]
    for title, body in sections:
        add_h1(doc, title)
        add_body(doc, body)

    add_h2(doc, "附：事实 / 判断 / 建议 标注示例")
    add_bullet(doc, "事实：抖音曝光 262,000 次，转化率 25.0%。")
    add_bullet(doc, "判断：抖音曝光高但有效线索转化偏低。")
    add_bullet(doc, "建议：优化内容与目标人群匹配，先做 A/B 测试，不急于加预算。")

    path = os.path.join(TPL_DIR, "10_市场分析报告模板.docx")
    set_core_props(doc, "市场分析报告模板")
    doc.save(path)
    return path


# ============================================================
# 11_月度复盘PPT模板.pptx
# ============================================================
def gen_tpl_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY_HEX = PptRGB(0x1B, 0x3A, 0x6B)
    BLUE_HEX = PptRGB(0x2E, 0x6B, 0xE6)
    GREY_HEX = PptRGB(0x66, 0x66, 0x66)

    slides = [
        ("封面", "青屿茶研市场部 7 月月度复盘　｜　（副标题请填写）"),
        ("目录", "（请填写：本复盘涵盖的章节）"),
        ("案例背景", "（请填写：品牌、部门、周期、目标）"),
        ("本月市场活动", "（请填写：5 项活动完成/延期情况）"),
        ("渠道数据概览", "（请填写：曝光、线索、有效线索、转化率，注明口径）"),
        ("已确认决定", "（请填写：逐条列出会议决定）"),
        ("行动事项", "（请填写：事项 / 负责人 / 截止，用表格）"),
        ("待确认事项", "（请填写：未决事项及原因）"),
        ("关键结论", "（请填写：每页一个核心结论）"),
        ("问题与挑战", "（请填写：数据反映的问题）"),
        ("下月计划", "（请填写：主要工作安排）"),
        ("结语", "（请填写：寄语或下一步）"),
    ]
    for title, body in slides:
        s = prs.slides.add_slide(blank)
        # 顶部色条
        bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY_HEX
        bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        p = tf.paragraphs[0]; p.text = title
        p.font.size = PptPt(28); p.font.bold = True
        p.font.color.rgb = PptRGB(0xFF, 0xFF, 0xFF)
        p.font.name = "微软雅黑"
        # 正文占位
        box = s.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                   Inches(12.1), Inches(5.5))
        tf2 = box.text_frame; tf2.word_wrap = True
        pp = tf2.paragraphs[0]; pp.text = body
        pp.font.size = PptPt(20); pp.font.color.rgb = GREY_HEX
        pp.font.name = "微软雅黑"
        # 角标
        tag = s.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6), Inches(0.4))
        tp = tag.text_frame.paragraphs[0]
        tp.text = "占位模板 · 由 WorkBuddy 填充"
        tp.font.size = PptPt(12); tp.font.color.rgb = BLUE_HEX
        tp.font.name = "微软雅黑"

    path = os.path.join(TPL_DIR, "11_月度复盘PPT模板.pptx")
    prs.save(path)
    return path


if __name__ == "__main__":
    C.verify()
    os.makedirs(DOC_DIR, exist_ok=True)
    os.makedirs(TPL_DIR, exist_ok=True)
    for fn in (gen_background, gen_meeting_raw, gen_work_log,
               gen_tpl_minutes, gen_tpl_summary, gen_tpl_notice,
               gen_tpl_report, gen_tpl_ppt):
        p = fn()
        print(f"[OK] {os.path.relpath(p, ROOT)}")
    print("\n文档与模板生成完成。")
